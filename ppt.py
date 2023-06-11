import pandas as pd
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

df = pd.read_csv('data.csv')

def createPPT(field):
    prs = Presentation()

    fieldDF = df[df['dep'] == field]
    print(fieldDF)

    for row in fieldDF.iterrows():
        quiz = row[1][0]
        dep = row[1][1]
        level = row[1][2]
        problemSet = row[1][3]

        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)

        title = slide.placeholders[0]

        tParagraph = title.text_frame.paragraphs[0]

        tParagraph.text = quiz
        tParagraph.font.name = 'Arial'
        tParagraph.font.size = Pt(40)
        tParagraph.font.bold = True
        tParagraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

        subtitle = slide.placeholders[1]

        sParagraph = subtitle.text_frame.paragraphs[0]

        sParagraph.text = f"분야 - {dep}\n난이도 - {level}\n{problemSet}"
        sParagraph.font.name = 'Arial'
        sParagraph.font.size = Pt(20)
        sParagraph.font.bold = False
        sParagraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    prs.save(f'{field}.pptx')

for field in df['dep'].unique():
    createPPT(field)